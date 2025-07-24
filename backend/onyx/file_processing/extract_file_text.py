import io
import json
import os
import re
import zipfile
from collections.abc import Callable
from collections.abc import Iterator
from collections.abc import Sequence
from email.parser import Parser as EmailParser
from enum import auto
from enum import IntFlag
from io import BytesIO
from pathlib import Path
from typing import Any
from typing import IO
from typing import NamedTuple
from zipfile import BadZipFile

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from docx import Document as DocxDocument
from fastapi import UploadFile
from PIL import Image
from pypdf import PdfReader
from pypdf.errors import PdfStreamError
import fitz  # PyMuPDF for better PDF handling and OCR fallback
import pytesseract  # OCR library

from onyx.configs.constants import FileOrigin
from onyx.configs.constants import ONYX_METADATA_FILENAME
from onyx.configs.llm_configs import get_image_extraction_and_analysis_enabled
from onyx.file_processing.html_utils import parse_html_page_basic
from onyx.file_processing.unstructured import get_unstructured_api_key
from onyx.file_processing.unstructured import unstructured_to_text
from onyx.file_store.file_store import FileStore
from onyx.utils.logger import setup_logger

logger = setup_logger()

# NOTE(rkuo): Unify this with upload_files_for_chat and file_valiation.py
TEXT_SECTION_SEPARATOR = "\n\n"

# OCR Configuration
OCR_CONFIG = {
    "MIN_TEXT_RATIO": 0.5,  # Minimum text per page ratio to consider text extraction successful  
    "MIN_WORDS_PER_PAGE": 10,  # Minimum words per page for successful text extraction
    "OCR_DPI": 300,  # DPI for OCR processing
    "OCR_LANG": "eng",  # Tesseract language
    "MAX_OCR_PAGES": 50,  # Maximum pages to process with OCR (to prevent excessive processing)
}

ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]

ACCEPTED_DOCUMENT_FILE_EXTENSIONS = [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
]

ACCEPTED_IMAGE_FILE_EXTENSIONS = [
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
]

ALL_ACCEPTED_FILE_EXTENSIONS = (
    ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS
    + ACCEPTED_DOCUMENT_FILE_EXTENSIONS
    + ACCEPTED_IMAGE_FILE_EXTENSIONS
)

IMAGE_MEDIA_TYPES = [
    "image/png",
    "image/jpeg",
    "image/webp",
]


class OnyxExtensionType(IntFlag):
    Plain = auto()
    Document = auto()
    Multimedia = auto()
    All = Plain | Document | Multimedia


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension.lower()


def is_valid_media_type(media_type: str) -> bool:
    return media_type in IMAGE_MEDIA_TYPES


def is_accepted_file_ext(ext: str, ext_type: OnyxExtensionType) -> bool:
    if ext_type & OnyxExtensionType.Plain:
        if ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Document:
        if ext in ACCEPTED_DOCUMENT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Multimedia:
        if ext in ACCEPTED_IMAGE_FILE_EXTENSIONS:
            return True

    return False


def is_text_file(file: IO[bytes]) -> bool:
    """
    checks if the first 1024 bytes only contain printable or whitespace characters
    if it does, then we say it's a plaintext file
    """
    raw_data = file.read(1024)
    file.seek(0)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    file.seek(0)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any]]]:
    """
    Iterates through files in a zip archive, yielding (ZipInfo, file handle) pairs.
    """
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        for file_info in zip_file.infolist():
            if ignore_dirs and file_info.is_dir():
                continue

            if (
                ignore_macos_resource_fork_files
                and is_macos_resource_fork_file(file_info.filename)
            ) or file_info.filename == ONYX_METADATA_FILENAME:
                continue

            with zip_file.open(file_info.filename, "r") as subfile:
                # Try to match by exact filename first
                yield file_info, subfile


def _extract_onyx_metadata(line: str) -> dict | None:
    """
    Example: first line has:
        <!-- ONYX_METADATA={"title": "..."} -->
      or
        #ONYX_METADATA={"title":"..."}
    """
    html_comment_pattern = r"<!--\s*ONYX_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#ONYX_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_onyx_metadata: bool = True,
) -> tuple[str, dict]:
    """
    For plain text files. Optionally extracts Onyx metadata from the first line.
    """
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        # decode
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        # optionally parse metadata in the first line
        if ind == 0 and not ignore_onyx_metadata:
            potential_meta = _extract_onyx_metadata(line)
            if potential_meta is not None:
                metadata = potential_meta
                continue

        file_content_raw += line

    return file_content_raw, metadata


def _should_use_ocr_fallback(text_content: str, total_pages: int) -> bool:
    """
    Determine if OCR fallback should be used based on extracted text quality.
    
    Args:
        text_content: Extracted text from PDF
        total_pages: Total number of pages in the PDF
        
    Returns:
        True if OCR fallback should be used
    """
    if not text_content.strip():
        logger.info("No text extracted, using OCR fallback")
        return True
    
    words = text_content.split()
    total_words = len(words)
    words_per_page = total_words / total_pages if total_pages > 0 else 0
    
    # Check if we have enough words per page
    if words_per_page < OCR_CONFIG["MIN_WORDS_PER_PAGE"]:
        logger.info(f"Low word density ({words_per_page:.1f} words/page), using OCR fallback")
        return True
    
    # Check text ratio (non-whitespace characters)
    non_whitespace_chars = len(re.sub(r'\s', '', text_content))
    total_chars = len(text_content)
    text_ratio = non_whitespace_chars / total_chars if total_chars > 0 else 0
    
    if text_ratio < OCR_CONFIG["MIN_TEXT_RATIO"]:
        logger.info(f"Low text ratio ({text_ratio:.2f}), using OCR fallback")
        return True
    
    return False


def _extract_text_with_ocr(pdf_document: fitz.Document, file_name: str = "") -> str:
    """
    Extract text from PDF using OCR on each page.
    
    Args:
        pdf_document: PyMuPDF document object
        file_name: Name of the file for logging
        
    Returns:
        Extracted text from all pages
    """
    logger.info(f"Starting OCR processing for {file_name}")
    
    if len(pdf_document) > OCR_CONFIG["MAX_OCR_PAGES"]:
        logger.warning(
            f"PDF has {len(pdf_document)} pages, limiting OCR to first {OCR_CONFIG['MAX_OCR_PAGES']} pages"
        )
        pages_to_process = OCR_CONFIG["MAX_OCR_PAGES"]
    else:
        pages_to_process = len(pdf_document)
    
    ocr_text_parts = []
    
    for page_num in range(pages_to_process):
        try:
            page = pdf_document[page_num]
            
            # Convert page to image
            mat = fitz.Matrix(OCR_CONFIG["OCR_DPI"] / 72, OCR_CONFIG["OCR_DPI"] / 72)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Convert to PIL Image for OCR
            image = Image.open(io.BytesIO(img_data))
            
            # Perform OCR
            page_text = pytesseract.image_to_string(
                image, 
                lang=OCR_CONFIG["OCR_LANG"],
                config='--psm 6'  # Assume uniform block of text
            )
            
            if page_text.strip():
                ocr_text_parts.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}")
                logger.debug(f"OCR extracted {len(page_text.split())} words from page {page_num + 1}")
            else:
                logger.debug(f"No text extracted via OCR from page {page_num + 1}")
                
        except Exception as e:
            logger.error(f"OCR failed for page {page_num + 1} in {file_name}: {e}")
            # Continue with other pages
            continue
    
    if pages_to_process < len(pdf_document):
        ocr_text_parts.append(f"\n--- Note: Only processed first {pages_to_process} pages of {len(pdf_document)} total pages ---")
    
    final_text = TEXT_SECTION_SEPARATOR.join(ocr_text_parts)
    logger.info(f"OCR completed for {file_name}. Extracted {len(final_text.split())} total words from {len(ocr_text_parts)} pages")
    
    return final_text


def _extract_images_from_pdf_pymupdf(pdf_document: fitz.Document) -> list[tuple[bytes, str]]:
    """
    Extract images from PDF using PyMuPDF.
    
    Args:
        pdf_document: PyMuPDF document object
        
    Returns:
        List of (image_bytes, image_name) tuples
    """
    extracted_images = []
    
    for page_num in range(len(pdf_document)):
        try:
            page = pdf_document[page_num]
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                xref = img[0]  # xref number
                pix = fitz.Pixmap(pdf_document, xref)
                
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("png")
                    image_name = f"page_{page_num + 1}_image_{img_index + 1}.png"
                    extracted_images.append((img_data, image_name))
                else:  # CMYK: convert to RGB first
                    pix1 = fitz.Pixmap(fitz.csRGB, pix)
                    img_data = pix1.tobytes("png")
                    image_name = f"page_{page_num + 1}_image_{img_index + 1}.png"
                    extracted_images.append((img_data, image_name))
                    pix1 = None
                    
                pix = None
                
        except Exception as e:
            logger.error(f"Failed to extract images from page {page_num + 1}: {e}")
            continue
    
    return extracted_images


def read_pdf_file_enhanced(
    file: IO[Any], 
    pdf_pass: str | None = None, 
    extract_images: bool = False,
    file_name: str = ""
) -> tuple[str, dict[str, Any], Sequence[tuple[bytes, str]]]:
    """
    Enhanced PDF reader with OCR fallback for scanned/blurry PDFs.
    
    Args:
        file: PDF file object
        pdf_pass: Password for encrypted PDFs
        extract_images: Whether to extract embedded images
        file_name: Name of the file for logging
        
    Returns:
        Tuple of (text_content, metadata, extracted_images)
    """
    metadata: dict[str, Any] = {}
    extracted_images: list[tuple[bytes, str]] = []
    
    try:
        # First try with pypdf for text extraction
        file.seek(0)
        pdf_reader = PdfReader(file)
        
        # Handle encryption
        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            try:
                decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
            except Exception:
                logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                return "", metadata, []
        elif pdf_reader.is_encrypted:
            logger.warning("No Password for an encrypted PDF, returning empty text.")
            return "", metadata, []

        # Extract basic PDF metadata
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value
                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        # Extract text using pypdf
        text_pages = []
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            text_pages.append(page_text)
        
        initial_text = TEXT_SECTION_SEPARATOR.join(text_pages)
        total_pages = len(pdf_reader.pages)
        
        # Check if we should use OCR fallback
        use_ocr = _should_use_ocr_fallback(initial_text, total_pages)
        
        final_text = initial_text
        
        if use_ocr:
            try:
                # Use PyMuPDF for OCR processing
                file.seek(0)
                file_bytes = file.read()
                pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
                
                # Extract text with OCR
                ocr_text = _extract_text_with_ocr(pdf_document, file_name)
                
                if ocr_text.strip():
                    final_text = ocr_text
                    metadata["processing_method"] = "OCR"
                    logger.info(f"Successfully extracted text via OCR from {file_name}")
                else:
                    logger.warning(f"OCR did not extract any text from {file_name}, using original text")
                    metadata["processing_method"] = "text_extraction_with_ocr_attempted"
                
                # Extract images if requested using PyMuPDF (better than pypdf for images)
                if extract_images:
                    try:
                        extracted_images = _extract_images_from_pdf_pymupdf(pdf_document)
                        logger.info(f"Extracted {len(extracted_images)} images from {file_name}")
                    except Exception as e:
                        logger.error(f"Failed to extract images with PyMuPDF from {file_name}: {e}")
                
                pdf_document.close()
                
            except Exception as e:
                logger.error(f"OCR processing failed for {file_name}: {e}")
                metadata["processing_method"] = "text_extraction_only"
                metadata["ocr_error"] = str(e)
        else:
            metadata["processing_method"] = "text_extraction"
            
            # Extract images using pypdf if requested and OCR wasn't used
            if extract_images:
                try:
                    for page_num, page in enumerate(pdf_reader.pages):
                        for image_file_object in page.images:
                            image = Image.open(io.BytesIO(image_file_object.data))
                            img_byte_arr = io.BytesIO()
                            image.save(img_byte_arr, format=image.format)
                            img_bytes = img_byte_arr.getvalue()

                            image_name = (
                                f"page_{page_num + 1}_image_{image_file_object.name}."
                                f"{image.format.lower() if image.format else 'png'}"
                            )
                            extracted_images.append((img_bytes, image_name))
                except Exception as e:
                    logger.error(f"Failed to extract images with pypdf from {file_name}: {e}")

        return final_text, metadata, extracted_images

    except PdfStreamError:
        logger.exception(f"Invalid PDF file: {file_name}")
    except Exception:
        logger.exception(f"Failed to read PDF: {file_name}")

    return "", metadata, []


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None, file_name: str = "") -> str:
    """
    Extract text from a PDF with OCR fallback for scanned documents.
    """
    text, _, _ = read_pdf_file_enhanced(file, pdf_pass, extract_images=False, file_name=file_name)
    return text


def docx_to_text_and_images(
    file: IO[Any], file_name: str = ""
) -> tuple[str, Sequence[tuple[bytes, str]]]:
    """
    Extract text from a docx. If embed_images=True, also extract inline images.
    Return (text_content, list_of_images).
    """
    paragraphs = []
    embedded_images: list[tuple[bytes, str]] = []

    try:
        doc = docx.Document(file)
    except BadZipFile as e:
        logger.warning(
            f"Failed to extract docx {file_name or 'docx file'}: {e}. Attempting to read as text file."
        )

        # May be an invalid docx, but still a valid text file
        file.seek(0)
        encoding = detect_encoding(file)
        text_content_raw, _ = read_text_file(
            file, encoding=encoding, ignore_onyx_metadata=False
        )
        return text_content_raw or "", []

    # Grab text from paragraphs
    for paragraph in doc.paragraphs:
        paragraphs.append(paragraph.text)

    # Reset position so we can re-load the doc (python-docx has read the stream)
    # Note: if python-docx has fully consumed the stream, you may need to open it again from memory.
    # For large docs, a more robust approach is needed.
    # This is a simplified example.

    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            # image is typically in rel.target_part.blob
            image_bytes = rel.target_part.blob
            image_name = rel.target_part.partname
            # store
            embedded_images.append((image_bytes, os.path.basename(str(image_name))))

    text_content = "\n".join(paragraphs)
    return text_content, embedded_images


def pptx_to_text(file: IO[Any], file_name: str = "") -> str:
    try:
        presentation = pptx.Presentation(file)
    except BadZipFile as e:
        error_str = f"Failed to extract text from {file_name or 'pptx file'}: {e}"
        logger.warning(error_str)
        return ""
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_content.append(slide_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def xlsx_to_text(file: IO[Any], file_name: str = "") -> str:
    try:
        workbook = openpyxl.load_workbook(file, read_only=True)
    except BadZipFile as e:
        error_str = f"Failed to extract text from {file_name or 'xlsx file'}: {e}"
        if file_name.startswith("~"):
            logger.debug(error_str + " (this is expected for files with ~)")
        else:
            logger.warning(error_str)
        return ""
    except Exception as e:
        if "File contains no valid workbook part" in str(e):
            logger.error(
                f"Failed to extract text from {file_name or 'xlsx file'}. This happens due to a bug in openpyxl. {e}"
            )
            return ""
        raise e

    text_content = []
    for sheet in workbook.worksheets:
        rows = []
        num_empty_consecutive_rows = 0
        for row in sheet.iter_rows(min_row=1, values_only=True):
            row_str = ",".join(str(cell or "") for cell in row)

            # Only add the row if there are any values in the cells
            if len(row_str) >= len(row):
                rows.append(row_str)
                num_empty_consecutive_rows = 0
            else:
                num_empty_consecutive_rows += 1

            if num_empty_consecutive_rows > 100:
                # handle massive excel sheets with mostly empty cells
                logger.warning(
                    f"Found {num_empty_consecutive_rows} empty rows in {file_name},"
                    " skipping rest of file"
                )
                break
        sheet_str = "\n".join(rows)
        text_content.append(sheet_str)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def eml_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    text_file = io.TextIOWrapper(file, encoding=encoding)
    parser = EmailParser()
    message = parser.parse(text_file)

    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            payload = part.get_payload()
            if isinstance(payload, str):
                text_content.append(payload)
            elif isinstance(payload, list):
                text_content.extend(item for item in payload if isinstance(item, str))
            else:
                logger.warning(f"Unexpected payload type: {type(payload)}")
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content, _ = read_text_file(file, encoding=encoding)
    return file_content


def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    """
    Legacy function that returns *only text*, ignoring embedded images.
    Now uses enhanced PDF processing with OCR fallback.

    NOTE: Ignoring seems to be defined as returning an empty string for files it can't
    handle (such as images).
    """
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": lambda f: pdf_to_text(f, file_name=file_name),  # Enhanced with OCR
        ".docx": lambda f: docx_to_text_and_images(f, file_name)[0],  # no images
        ".pptx": lambda f: pptx_to_text(f, file_name),
        ".xlsx": lambda f: xlsx_to_text(f, file_name),
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
    }

    try:
        if get_unstructured_api_key():
            try:
                return unstructured_to_text(file, file_name)
            except Exception as unstructured_error:
                logger.error(
                    f"Failed to process with Unstructured: {str(unstructured_error)}. "
                    "Falling back to normal processing."
                )
        if extension is None:
            extension = get_file_ext(file_name)

        if is_accepted_file_ext(
            extension, OnyxExtensionType.Plain | OnyxExtensionType.Document
        ):
            func = extension_to_function.get(extension, file_io_to_text)
            file.seek(0)
            return func(file)

        # If unknown extension, maybe it's a text file
        file.seek(0)
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension or not recognized as text data")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""


class ExtractionResult(NamedTuple):
    """Structured result from text and image extraction from various file types."""

    text_content: str
    embedded_images: Sequence[tuple[bytes, str]]
    metadata: dict[str, Any]


def extract_text_and_images(
    file: IO[Any],
    file_name: str,
    pdf_pass: str | None = None,
) -> ExtractionResult:
    """
    Primary new function for the updated connector with enhanced PDF processing.
    Returns structured extraction result with text content, embedded images, and metadata.
    """

    try:
        # Attempt unstructured if env var is set
        if get_unstructured_api_key():
            # If the user doesn't want embedded images, unstructured is fine
            file.seek(0)
            text_content = unstructured_to_text(file, file_name)
            return ExtractionResult(
                text_content=text_content, embedded_images=[], metadata={}
            )

        extension = get_file_ext(file_name)

        # docx example for embedded images
        if extension == ".docx":
            file.seek(0)
            text_content, images = docx_to_text_and_images(file, file_name)
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata={}
            )

        # Enhanced PDF processing with OCR fallback
        if extension == ".pdf":
            file.seek(0)
            text_content, pdf_metadata, images = read_pdf_file_enhanced(
                file,
                pdf_pass,
                extract_images=get_image_extraction_and_analysis_enabled(),
                file_name=file_name
            )
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata=pdf_metadata
            )

        # For PPTX, XLSX, EML, etc., we do not show embedded image logic here.
        # You can do something similar to docx if needed.
        if extension == ".pptx":
            file.seek(0)
            return ExtractionResult(
                text_content=pptx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".xlsx":
            file.seek(0)
            return ExtractionResult(
                text_content=xlsx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".eml":
            file.seek(0)
            return ExtractionResult(
                text_content=eml_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".epub":
            file.seek(0)
            return ExtractionResult(
                text_content=epub_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".html":
            file.seek(0)
            return ExtractionResult(
                text_content=parse_html_page_basic(file),
                embedded_images=[],
                metadata={},
            )

        # If we reach here and it's a recognized text extension
        if is_text_file_extension(file_name):
            file.seek(0)
            encoding = detect_encoding(file)
            text_content_raw, file_metadata = read_text_file(
                file, encoding=encoding, ignore_onyx_metadata=False
            )
            return ExtractionResult(
                text_content=text_content_raw,
                embedded_images=[],
                metadata=file_metadata,
            )

        # If it's an image file or something else, we do not parse embedded images from them
        # just return empty text
        file.seek(0)
        return ExtractionResult(text_content="", embedded_images=[], metadata={})

    except Exception as e:
        logger.exception(f"Failed to extract text/images from {file_name}: {e}")
        return ExtractionResult(text_content="", embedded_images=[], metadata={})


def convert_docx_to_txt(file: UploadFile, file_store: FileStore) -> str:
    """
    Helper to convert docx to a .txt file in the same filestore.
    """
    file.file.seek(0)
    docx_content = file.file.read()
    doc = DocxDocument(BytesIO(docx_content))

    # Extract text from the document
    all_paras = [p.text for p in doc.paragraphs]
    text_content = "\n".join(all_paras)

    file_id = file_store.save_file(
        content=BytesIO(text_content.encode("utf-8")),
        display_name=file.filename,
        file_origin=FileOrigin.CONNECTOR,
        file_type="text/plain",
    )
    return file_id


def docx_to_txt_filename(file_path: str) -> str:
    return file_path.rsplit(".", 1)[0] + ".txt"
